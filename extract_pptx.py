import sys
from pptx import Presentation

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = f"Slide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
        text.append(slide_text)
    return "\n".join(text)

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else '/vercel/sandbox/uploads/C_Store_AI_System_Presentation 2.pptx'
    print(extract_text(pptx_path))